#!/usr/bin/env python3
"""Build CipherOutreach COPD & CHF patient journey slides (CipherHealth 2026 brand)."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "CipherOutreach_DRG_Patient_Journeys.pptx"

INDIGO = RGBColor(0x2C, 0x26, 0x60)
VIOLET = RGBColor(0x5B, 0x4F, 0xB5)
TEAL = RGBColor(0x00, 0xB3, 0xA3)
LIME = RGBColor(0xB5, 0xCC, 0x2E)
GOLD = RGBColor(0xE8, 0xA0, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL_TINT = RGBColor(0xBC, 0xFF, 0xF9)
CARD_LINE = RGBColor(0xE4, 0xE2, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COPD = {
    "program": "COPD",
    "overline": "CipherOutreach  ·  DRG-Specific Longitudinal Workflow",
    "title": "Patti's COPD Recovery Journey",
    "subtitle": (
        "Four automated check-ins over 23 days after hospital discharge — "
        "monitoring breathing, medications, and follow-up so the care team can intervene early."
    ),
    "duration": "23 days",
    "touches": "4 touches",
    "stops": [
        {
            "day": "Day 0",
            "name": "Hospital Discharge",
            "desc": "Patti leaves the hospital and is enrolled in CipherOutreach COPD recovery outreach.",
            "tone": "start",
        },
        {
            "day": "Day 2",
            "name": "Initial Recovery Check-In",
            "desc": "Asks how she's feeling, care-instruction understanding, smoking status, medications, shortness of breath, and follow-up scheduling.",
            "tone": "touch",
        },
        {
            "day": "Day 9",
            "name": "Early Progress Check-In",
            "desc": "Checks improvement since discharge, breathing progress, medication adherence, and whether her follow-up visit is complete.",
            "tone": "touch",
        },
        {
            "day": "Day 16",
            "name": "Continued Recovery Check-In",
            "desc": "Reassesses how she feels, shortness of breath over the last week, and daily medication adherence — plus a quit-smoking resource.",
            "tone": "touch",
        },
        {
            "day": "Day 23",
            "name": "Final Recovery Check-In",
            "desc": "Closing check on how she feels versus discharge, breathing, and medication adherence — then thanks her for completing the series.",
            "tone": "end",
        },
    ],
}

CHF = {
    "program": "CHF",
    "overline": "CipherOutreach  ·  DRG-Specific Longitudinal Workflow",
    "title": "Patti's CHF Recovery Journey",
    "subtitle": (
        "Five automated check-ins over 30 days after hospital discharge — "
        "tracking weight, breathing, diet, medications, and warning symptoms with care-team escalation when needed."
    ),
    "duration": "30 days",
    "touches": "5 touches",
    "stops": [
        {
            "day": "Day 0",
            "name": "Hospital Discharge",
            "desc": "Patti leaves the hospital and is enrolled in CipherOutreach CHF post-discharge recovery outreach.",
            "tone": "start",
        },
        {
            "day": "Day 2",
            "name": "Initial CHF Recovery Check-In",
            "desc": "Asks about how she feels, rapid weight gain, SOB, follow-up, low-salt diet, medications, care instructions, home health, and warning symptoms.",
            "tone": "touch",
        },
        {
            "day": "Day 7",
            "name": "Early Progress Check-In",
            "desc": "Rechecks daily weight trend, activity-limiting SOB, medication adherence, low-salt diet, home health, follow-up status, and CHF symptoms.",
            "tone": "touch",
        },
        {
            "day": "Day 14",
            "name": "Continued Recovery Check-In",
            "desc": "Reinforces daily weigh-ins, breathing, medication adherence, low-salt diet, and any swelling, chest pain, coughing, or tiredness.",
            "tone": "touch",
        },
        {
            "day": "Day 21",
            "name": "Late Recovery Check-In",
            "desc": "Continues monitoring weight, shortness of breath, medications, diet adherence, and CHF warning symptoms as recovery progresses.",
            "tone": "touch",
        },
        {
            "day": "Day 30",
            "name": "Final Recovery Check-In",
            "desc": "Final check on weight stability, breathing, medications, diet, and symptoms — then closes the 30-day recovery journey.",
            "tone": "end",
        },
    ],
}


def set_run_font(run, name: str, size_pt: float, bold=False, italic=False, color=INDIGO):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs", "a:ea"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font="Inter",
    size=12,
    bold=False,
    italic=False,
    color=INDIGO,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    # Tighten margins
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    run = p.add_run()
    run.text = text
    set_run_font(run, font, size, bold=bold, italic=italic, color=color)
    return box


def add_title_with_accent(slide, left, top, width, height, before, accent, after=""):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if before:
        r = p.add_run()
        r.text = before
        set_run_font(r, "Merriweather", 26, bold=True, color=INDIGO)
    r = p.add_run()
    r.text = accent
    set_run_font(r, "Merriweather", 26, bold=True, italic=True, color=VIOLET)
    if after:
        r = p.add_run()
        r.text = after
        set_run_font(r, "Merriweather", 26, bold=True, color=INDIGO)
    return box


def rounded_rect(slide, left, top, width, height, fill, line=None, adj=0.12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    try:
        shape.adjustments[0] = adj
    except Exception:
        pass
    return shape


def oval(slide, left, top, width, height, fill, line=None, line_width=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    return shape


def tone_colors(tone: str):
    """Return (accent_fill, chip_fill, chip_text, node_fill, node_text)."""
    if tone == "start":
        # Gold chip/node with indigo ink (contrast-safe)
        return GOLD, GOLD, INDIGO, GOLD, INDIGO
    if tone == "end":
        return LIME, LIME, INDIGO, LIME, INDIGO
    # touch: violet chips (white text OK), teal path accents
    return TEAL, VIOLET, WHITE, VIOLET, WHITE


def build_slide(prs: Presentation, journey: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background (full-bleed brand motif)
    slide.shapes.add_picture(str(ROOT / "bg_content.png"), 0, 0, SLIDE_W, SLIDE_H)

    # Logo
    slide.shapes.add_picture(str(ROOT / "logo.png"), Inches(11.45), Inches(0.26), Inches(1.55))

    # Overline — violet on lavender (teal fails contrast on lavender)
    add_textbox(
        slide,
        Inches(0.5),
        Inches(0.26),
        Inches(9.5),
        Inches(0.28),
        journey["overline"],
        font="Inter",
        size=11,
        bold=True,
        color=VIOLET,
    )

    prog = journey["program"]
    title = journey["title"]
    if prog in title:
        before, after = title.split(prog, 1)
        add_title_with_accent(
            slide,
            Inches(0.5),
            Inches(0.52),
            Inches(10.6),
            Inches(0.5),
            before,
            prog,
            after,
        )
    else:
        add_textbox(
            slide,
            Inches(0.5),
            Inches(0.52),
            Inches(10.6),
            Inches(0.5),
            title,
            font="Merriweather",
            size=26,
            bold=True,
            color=INDIGO,
        )

    add_textbox(
        slide,
        Inches(0.5),
        Inches(1.05),
        Inches(8.7),
        Inches(0.55),
        journey["subtitle"],
        font="Inter",
        size=12,
        color=INDIGO,
    )

    # Meta chips — violet / indigo fills with white text (contrast-safe)
    for i, (label, fill) in enumerate(
        [(journey["duration"], VIOLET), (journey["touches"], INDIGO)]
    ):
        left = Inches(9.45)
        top = Inches(1.08) + Inches(i * 0.38)
        rounded_rect(slide, left, top, Inches(1.55), Inches(0.32), fill)
        add_textbox(
            slide,
            left,
            top,
            Inches(1.55),
            Inches(0.32),
            label,
            font="Inter",
            size=11,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Patient card
    rounded_rect(
        slide,
        Inches(0.4),
        Inches(1.75),
        Inches(2.2),
        Inches(5.0),
        WHITE,
        line=CARD_LINE,
        adj=0.08,
    )
    # Teal accent strip
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.75), Inches(0.1), Inches(5.0)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = TEAL
    strip.line.fill.background()

    av_size = Inches(1.5)
    slide.shapes.add_picture(
        str(ROOT / "patti_avatar.png"), Inches(0.75), Inches(2.0), av_size, av_size
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(3.65),
        Inches(2.0),
        Inches(0.32),
        "Meet Patti",
        font="Merriweather",
        size=15,
        bold=True,
        color=INDIGO,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(3.98),
        Inches(2.0),
        Inches(0.26),
        "Patient journey avatar",
        font="Inter",
        size=10,
        bold=True,
        color=VIOLET,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.55),
        Inches(4.35),
        Inches(1.9),
        Inches(1.6),
        "Follow Patti from discharge through each CipherOutreach touchpoint — "
        "what we ask of her at every stop along the path.",
        font="Inter",
        size=11,
        color=INDIGO,
        align=PP_ALIGN.CENTER,
    )
    # Soft teal tint tip
    tip = rounded_rect(
        slide, Inches(0.55), Inches(6.15), Inches(1.9), Inches(0.4), TEAL_TINT, adj=0.2
    )
    add_textbox(
        slide,
        Inches(0.55),
        Inches(6.15),
        Inches(1.9),
        Inches(0.4),
        "Linear recovery path",
        font="Inter",
        size=10,
        bold=True,
        color=INDIGO,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    stops = journey["stops"]
    n = len(stops)

    # Timeline band
    path_left = Inches(2.95)
    path_right = Inches(12.9)
    path_width = path_right - path_left
    path_y = Inches(2.55)

    # Soft track underlay
    under = rounded_rect(
        slide,
        path_left - Inches(0.08),
        path_y - Inches(0.22),
        path_width + Inches(0.16),
        Inches(0.44),
        TEAL_TINT,
        adj=0.5,
    )

    track = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        path_left,
        path_y - Inches(0.035),
        path_width,
        Inches(0.07),
    )
    track.fill.solid()
    track.fill.fore_color.rgb = TEAL
    track.line.fill.background()

    usable = path_width - Inches(0.15)
    step = usable / max(n - 1, 1)

    # Card sizing
    gutter = Inches(0.12)
    area_left = Inches(2.85)
    area_width = Inches(10.15)
    card_w = (area_width - gutter * (n - 1)) / n
    card_h = Inches(3.55)
    card_top = Inches(3.15)

    for i, stop in enumerate(stops):
        accent, chip_fill, chip_text, node_fill, node_text = tone_colors(stop["tone"])
        cx = path_left + Inches(0.075) + step * i

        # Node
        node_r = Inches(0.2)
        oval(
            slide,
            cx - node_r - Inches(0.04),
            path_y - node_r - Inches(0.04),
            (node_r + Inches(0.04)) * 2,
            (node_r + Inches(0.04)) * 2,
            WHITE,
            line=accent,
            line_width=Pt(2.25),
        )
        oval(slide, cx - node_r, path_y - node_r, node_r * 2, node_r * 2, node_fill)
        add_textbox(
            slide,
            cx - node_r,
            path_y - node_r,
            node_r * 2,
            node_r * 2,
            str(i + 1),
            font="Inter",
            size=11,
            bold=True,
            color=node_text,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Connector down to card
        stub = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            cx - Inches(0.012),
            path_y + Inches(0.2),
            Inches(0.024),
            Inches(0.4),
        )
        stub.fill.solid()
        stub.fill.fore_color.rgb = accent
        stub.line.fill.background()

        card_left = area_left + i * (card_w + gutter)
        rounded_rect(slide, card_left, card_top, card_w, card_h, WHITE, line=CARD_LINE, adj=0.08)

        # Top accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, card_left, card_top, card_w, Inches(0.09)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # Day chip
        chip_w = Inches(0.9)
        chip_left = card_left + (card_w - chip_w) / 2
        rounded_rect(
            slide,
            chip_left,
            card_top + Inches(0.22),
            chip_w,
            Inches(0.28),
            chip_fill,
            adj=0.35,
        )
        add_textbox(
            slide,
            chip_left,
            card_top + Inches(0.22),
            chip_w,
            Inches(0.28),
            stop["day"],
            font="Inter",
            size=10,
            bold=True,
            color=chip_text,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            card_left + Inches(0.1),
            card_top + Inches(0.6),
            card_w - Inches(0.2),
            Inches(0.7),
            stop["name"],
            font="Merriweather",
            size=12 if n <= 5 else 11,
            bold=True,
            color=INDIGO,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            card_left + Inches(0.12),
            card_top + Inches(1.35),
            card_w - Inches(0.24),
            card_h - Inches(1.5),
            stop["desc"],
            font="Inter",
            size=10 if n <= 5 else 9,
            color=INDIGO,
            align=PP_ALIGN.CENTER,
        )

    # Footer
    add_textbox(
        slide,
        Inches(0.5),
        Inches(7.05),
        Inches(4),
        Inches(0.25),
        "Patient journey  ·  CipherOutreach touchpoints",
        font="Inter",
        size=9,
        color=VIOLET,
    )
    add_textbox(
        slide,
        Inches(4.2),
        Inches(7.05),
        Inches(5),
        Inches(0.25),
        "CIPHERHEALTH CONFIDENTIAL",
        font="Inter",
        size=9,
        bold=True,
        color=INDIGO,
        align=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build_slide(prs, COPD)
    build_slide(prs, CHF)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
